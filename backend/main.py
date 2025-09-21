from fastapi import FastAPI, Depends, HTTPException, Query, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session
from sqlalchemy import select, distinct, delete
from database import Base, engine, get_db
from models import Student, Activity, ActivityParticipant
from schemas import ActivityOut, SaveParticipantsRequest, SaveParticipantsResponse, StudentOut, ParticipantState
from typing import Literal
from io import BytesIO
import re
import os
from copy import deepcopy
from datetime import date as dt_date
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


app = FastAPI(
    title="Activity Participation API",
    docs_url="/docs" if os.getenv("VERCEL_ENV") != "production" else None,
    redoc_url="/redoc" if os.getenv("VERCEL_ENV") != "production" else None,
)

# CORS (adjust origins for your frontend)
app.add_middleware(
    CORSMiddleware,
    allow_origin_regex=r"https?://(localhost|127\.0\.0\.1):\d+|https://.*\.vercel\.app",
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Lazy table creation to avoid cold start delays
def ensure_tables():
    try:
        Base.metadata.create_all(bind=engine)
    except Exception as e:
        print(f"Table creation warning: {e}")

# Create tables on startup (avoid work at import-time in serverless)
@app.on_event("startup")
def _on_startup():
    # Skip table creation in serverless to reduce cold start time
    # Tables should already exist from migration
    pass


@app.get("/health")
def health_check():
    """Simple health check endpoint for Vercel"""
    return {"status": "ok", "service": "activity-participation-api"}


@app.get("/rooms", response_model=list[int])
def get_rooms(school_year: str | None = Query(None), db: Session = Depends(get_db)):
    q = select(distinct(Student.room))
    if school_year:
        q = q.where(Student.school_year == school_year)
    rooms = db.execute(q).scalars().all()
    return sorted([r for r in rooms if r is not None])


@app.get("/activities", response_model=list[ActivityOut])
def get_activities(db: Session = Depends(get_db)):
    acts = (
        db.execute(
            select(Activity)
            .where(Activity.show_in_ui == True)
            .order_by(Activity.activity_date.desc(), Activity.name)
        )
        .scalars()
        .all()
    )
    return acts

@app.get("/years", response_model=list[str])
def get_years(db: Session = Depends(get_db)):
    years = db.execute(select(distinct(Student.school_year))).scalars().all()
    # Filter out nulls and sort descending by academic year string
    years = [y for y in years if y]
    return sorted(years, reverse=True)


@app.get("/students", response_model=list[StudentOut])
def get_students(
    room: int = Query(...),
    school_year: str | None = Query(None),
    on_date: str | None = None,  # currently informational
    db: Session = Depends(get_db),
):
    q = select(Student).where(Student.room == room)
    if school_year:
        q = q.where(Student.school_year == school_year)
    students = db.execute(q.order_by(Student.first_name, Student.last_name)).scalars().all()
    return students


@app.get("/activities/{activity_id}/participants", response_model=list[ParticipantState])
def get_participants(
    activity_id: int,
    room: int = Query(...),
    school_year: str | None = Query(None),
    db: Session = Depends(get_db),
):
    """Fetch current participation state for a given activity and room."""
    q = (
        select(ActivityParticipant)
        .join(Student, Student.id == ActivityParticipant.student_id)
        .where(
            ActivityParticipant.activity_id == activity_id,
            Student.room == room,
        )
    )
    if school_year:
        q = q.where(Student.school_year == school_year)
    participants = db.execute(q).scalars().all()
    return participants


@app.post("/activities/{activity_id}/participants", response_model=SaveParticipantsResponse)
def save_participants(activity_id: int, req: SaveParticipantsRequest, db: Session = Depends(get_db)):
    # Validate activity
    act = db.get(Activity, activity_id)
    if not act:
        raise HTTPException(status_code=404, detail="Activity not found")

    upserted = 0
    deleted = 0

    for item in req.participants:
        # Determine if this row should be stored
        should_store = bool(item.participated or item.position)

        # Team validation: require team name only when a position is assigned
        if should_store and act.is_team and item.position is not None:
            if not item.team_name or not item.team_name.strip():
                raise HTTPException(status_code=422, detail=f"Team name required when saving a position for student_id={item.student_id} in a team activity.")

        existing = db.get(ActivityParticipant, {"activity_id": activity_id, "student_id": item.student_id})

        if should_store:
            if existing:
                existing.position = int(item.position) if item.position else None
                existing.team_name = item.team_name.strip() if item.team_name else None
            else:
                ap = ActivityParticipant(
                    activity_id=activity_id,
                    student_id=item.student_id,
                    position=int(item.position) if item.position else None,
                    team_name=item.team_name.strip() if item.team_name else None,
                )
                db.add(ap)
            upserted += 1
        else:
            if existing:
                db.delete(existing)
                deleted += 1

    db.commit()
    return SaveParticipantsResponse(upserted=upserted, deleted=deleted)


def _ordinal(n: int) -> str:
    if 10 <= n % 100 <= 20:
        suf = "th"
    else:
        suf = {1: "st", 2: "nd", 3: "rd"}.get(n % 10, "th")
    return f"{n}{suf}"


def _duplicate_slide(prs: Presentation, src_slide):
    # Create new slide with same layout
    slide_layout = src_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    # Remove any auto-created placeholders
    for shp in list(new_slide.shapes):
        new_slide.shapes._spTree.remove(shp._element)
    # Copy shapes from source, preserving pictures via public API
    for shp in src_slide.shapes:
        try:
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_bytes = shp.image.blob
                new_slide.shapes.add_picture(BytesIO(img_bytes), shp.left, shp.top, width=shp.width, height=shp.height)
            else:
                new_el = deepcopy(shp._element)
                new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
        except Exception:
            # Fallback: deep copy on any unexpected case
            new_el = deepcopy(shp._element)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    return new_slide


def _replace_placeholders(slide, mapping: dict[str, str]) -> None:
    def replace_in_string(text: str) -> str:
        new_text = text
        for k, v in mapping.items():
            aliases = {k}
            if "_" in k:
                aliases.add(k.replace("_", " "))
            else:
                aliases.add(k.replace(" ", "_"))
            # Include upper variants
            aliases |= {a.upper() for a in list(aliases)}
            # Build regex to match {[ alias ]} and { {alias} } with optional spaces.
            # We escape aliases for regex safety and join with alternation.
            alt = "|".join(sorted({re.escape(a) for a in aliases}, key=len, reverse=True))
            patterns = [
                rf"\{{\s*(?:{alt})\s*\}}",  # {ALIAS}
                rf"\[\s*(?:{alt})\s*\]",    # [ALIAS]
            ]
            for pat in patterns:
                new_text = re.sub(pat, v, new_text)
        return new_text

    def walk(shape):
        # Recurse groups
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            for ch in shape.shapes:
                walk(ch)
            return
        # Tables
        if getattr(shape, "has_table", False):
            tbl = shape.table
            for row in tbl.rows:
                for cell in row.cells:
                    for p in cell.text_frame.paragraphs:
                        # Run-level replacement (preserve formatting)
                        for run in p.runs:
                            run.text = replace_in_string(run.text)
                        # Fallback: if any placeholders remain across runs, replace paragraph text
                        joined = "".join(r.text for r in p.runs) if p.runs else p.text
                        replaced = replace_in_string(joined)
                        if replaced != joined:
                            p.text = replaced
            return
        # Text frames
        if hasattr(shape, "text_frame") and shape.has_text_frame:
            tf = shape.text_frame
            for p in tf.paragraphs:
                # Run-level replacement (preserve formatting)
                for run in p.runs:
                    run.text = replace_in_string(run.text)
                # Fallback: if any placeholders remain across runs, replace paragraph text
                joined = "".join(r.text for r in p.runs) if p.runs else p.text
                replaced = replace_in_string(joined)
                if replaced != joined:
                    p.text = replaced

    for shape in slide.shapes:
        walk(shape)


def _replace_literal(slide, search: str, replace: str) -> None:
    """Replace a literal substring across all text runs in the slide."""
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or not shape.has_text_frame:
            continue
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                if search in run.text:
                    run.text = run.text.replace(search, replace)


TemplateKind = Literal[
    'position_individual',
    'position_team',
    'participation',
]


@app.post("/activities/{activity_id}/certificates/ppt")
def generate_certificates_ppt(
    activity_id: int,
    room: int = Form(...),
    template_kind: TemplateKind = Form(...),
    award_date: str | None = Form(None),
    school_year: str | None = Form(None),
    template: UploadFile = File(...),
    db: Session = Depends(get_db),
):
    act = None
    # For position templates, we require and validate the selected activity
    needs_team = template_kind == 'position_team'
    if template_kind in ('position_team', 'position_individual'):
        act = db.get(Activity, activity_id)
        if not act:
            raise HTTPException(status_code=404, detail="Activity not found")
        if not act.show_in_ui:
            raise HTTPException(status_code=400, detail="Selected activity is not available for certificates.")
        if needs_team and not act.is_team:
            raise HTTPException(status_code=400, detail="Selected template expects a team activity, but activity is individual.")
        if (not needs_team) and act.is_team:
            raise HTTPException(status_code=400, detail="Selected template expects an individual activity, but activity is team-based.")

    if template_kind in ('position_team', 'position_individual'):
        # Load students and their row for THIS activity only
        q = (
            select(Student, ActivityParticipant)
            .outerjoin(
                ActivityParticipant,
                (ActivityParticipant.student_id == Student.id) & (ActivityParticipant.activity_id == activity_id),
            )
            .where(Student.room == room)
        )
        if school_year:
            q = q.where(Student.school_year == school_year)
        rows = db.execute(q.order_by(Student.first_name, Student.last_name)).all()

        recipients: list[tuple[Student, ActivityParticipant | None]] = []
        for s, ap in rows:
            if ap and ap.position is not None:
                recipients.append((s, ap))
        if not recipients:
            raise HTTPException(status_code=400, detail="No matching recipients (no positions found) for the given room and activity.")
    else:
        # Participation certs:
        # Ignore the activity_id and gather ALL activities in which each student participated
        # without a position, for the selected room. Aggregate events per student.
        q = (
            select(Student, Activity)
            .join(ActivityParticipant, ActivityParticipant.student_id == Student.id)
            .join(Activity, Activity.id == ActivityParticipant.activity_id)
            .where(Student.room == room)
            .where(ActivityParticipant.position.is_(None))
            .where(Activity.show_in_ui == True)
        )
        if school_year:
            q = q.where(Student.school_year == school_year)
        rows = db.execute(q.order_by(Student.first_name, Student.last_name, Activity.activity_date)).all()

        # Build mapping: student_id -> (Student, list[Activity])
        events_by_student: dict[int, tuple[Student, list[Activity]]] = {}
        for s, act_row in rows:
            entry = events_by_student.get(s.id)
            if not entry:
                events_by_student[s.id] = (s, [act_row])
            else:
                # Avoid duplicate activities for a student if any
                student, acts = entry
                if not any(a.id == act_row.id for a in acts):
                    acts.append(act_row)
                events_by_student[s.id] = (student, acts)

        # Recipients: only students with at least one non-positional participation
        recipients: list[tuple[Student, list[Activity]]] = [val for val in events_by_student.values() if val[1]]

        if not recipients:
            raise HTTPException(status_code=400, detail="No matching recipients for participation certificates in this room.")

    if not recipients:
        raise HTTPException(status_code=400, detail="No matching recipients for the given room and template kind.")

    # Read template file into Presentation
    content = template.file.read()
    if not content:
        raise HTTPException(status_code=400, detail="Empty template file.")
    try:
        prs = Presentation(BytesIO(content))
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Invalid PPTX template: {e}")

    if len(prs.slides) == 0:
        raise HTTPException(status_code=400, detail="Template must contain at least one slide with placeholders.")

    template_slide = prs.slides[0]

    # Build slides
    created_count = 0
    if template_kind in ('position_team', 'position_individual'):
        # Position certs: event date from this activity, position label as First/Second/Third
        date_str = (act.activity_date.isoformat() if isinstance(act.activity_date, dt_date) else str(act.activity_date))
        POS_LABEL = {1: "First", 2: "Second", 3: "Third"}
        for idx, (s, ap) in enumerate(recipients):
            slide = _duplicate_slide(prs, template_slide)

            mapping = {
                "name": f"{s.first_name} {s.last_name}",
                "event date": date_str,
                "event": act.name,
                "year": str(act.year),
                "room": str(s.room),
            }
            if ap and ap.position is not None:
                mapping["position"] = POS_LABEL.get(int(ap.position), str(ap.position))
            if needs_team:
                mapping["team_name"] = ap.team_name if ap and ap.team_name else ""

            _replace_placeholders(slide, mapping)
            created_count += 1
    else:
        # Participation certs: events are aggregated per student across all non-positional participations.
        def _format_events_str(acts: list[Activity]) -> str:
            # Keep order by date (already ordered), dedupe names while preserving order
            seen: set[str] = set()
            names: list[str] = []
            for a in acts:
                nm = a.name
                if nm not in seen:
                    seen.add(nm)
                    names.append(nm)
            return ", ".join(names)

        for s, acts in recipients:
            slide = _duplicate_slide(prs, template_slide)

            # Determine award date per student: provided award_date or latest activity date
            if award_date:
                date_str = award_date
            else:
                if acts:
                    latest = max((a.activity_date for a in acts))
                    date_str = (latest.isoformat() if isinstance(latest, dt_date) else str(latest))
                else:
                    date_str = ""

            mapping = {
                "name": f"{s.first_name} {s.last_name}",
                "award_date": date_str,
                "events": _format_events_str(acts),
                "room": str(s.room),
            }
            _replace_placeholders(slide, mapping)
            # If the template has the literal word "Tournaments", singularize it to "Tournament"
            _replace_literal(slide, "Tournaments", "Tournament")
            created_count += 1

    # Remove the original template slide if we created any slides
    if created_count > 0 and len(prs.slides) > 0:
        sldIdLst = prs.slides._sldIdLst
        first = sldIdLst[0]
        sldIdLst.remove(first)

    # Stream result
    out = BytesIO()
    prs.save(out)
    out.seek(0)
    filename = f"certificates_room{room}_activity{activity_id}_{template_kind}.pptx"
    return StreamingResponse(out, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers={
        "Content-Disposition": f"attachment; filename=\"{filename}\""
    })
